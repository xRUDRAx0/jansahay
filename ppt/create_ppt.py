import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Define colors based on template
    bg_color = RGBColor(10, 20, 35) # Dark blue/navy background
    text_color = RGBColor(255, 255, 255) # White text
    accent_color = RGBColor(0, 255, 255) # Cyan accent
    gold_color = RGBColor(255, 215, 0) # Gold accent
    
    def apply_template_style(slide, title_text, is_title_slide=False):
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        if is_title_slide:
            return
            
        # Add title box (mimicking the pill-shaped title in the template)
        left = Inches(2.5)
        top = Inches(0.5)
        width = Inches(5)
        height = Inches(0.8)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(20, 30, 50)
        shape.line.color.rgb = gold_color
        shape.line.width = Pt(1.5)
        
        text_frame = shape.text_frame
        text_frame.text = title_text
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.size = Pt(28)
        text_frame.paragraphs[0].font.color.rgb = text_color
        
        # Add main content bounding box (mimicking the large cyan frame)
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9)
        height = Inches(5)
        
        content_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        content_box.fill.background()
        content_box.line.color.rgb = accent_color
        content_box.line.width = Pt(1)

    # 1. Title Slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_template_style(slide, "", is_title_slide=True)
    
    txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    tf = txBox.text_frame
    tf.text = "JANSAHAY"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.color.rgb = text_color
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "AI PUBLIC-SERVICE COPILOT"
    p.font.size = Pt(24)
    p.font.color.rgb = accent_color
    p.alignment = PP_ALIGN.CENTER
    
    txBox2 = slide.shapes.add_textbox(Inches(5.5), Inches(5.5), Inches(4), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.text = "TEAM NAME: [Your Team Name]"
    tf2.paragraphs[0].font.size = Pt(18)
    tf2.paragraphs[0].font.color.rgb = text_color
    p2 = tf2.add_paragraph()
    p2.text = "TEAM MEMBERS: [Your Members]"
    p2.font.size = Pt(18)
    p2.font.color.rgb = text_color

    def add_content_text(slide, content_lines):
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(4))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(content_lines):
            if i == 0:
                tf.text = line
                tf.paragraphs[0].font.size = Pt(20)
                tf.paragraphs[0].font.color.rgb = text_color
            else:
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(20)
                p.font.color.rgb = text_color
                p.space_before = Pt(14)

    # 2. PROBLEM STATEMENT
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_template_style(slide, "PROBLEM STATEMENT")
    content = [
        "• Citizens often struggle to navigate complex government services and eligibility criteria.",
        "• Manual document verification is slow, error-prone, and delays benefit distribution.",
        "• Lack of personalized guidance makes scheme discovery difficult for those in need.",
        "• Language barriers and low tech-literacy create friction in accessing public services."
    ]
    add_content_text(slide, content)

    # 3. SOLUTION
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_template_style(slide, "SOLUTION")
    content = [
        "• AI-Powered Copilot: Understands user needs in natural language and matches them with relevant services.",
        "• Document Doctor: Automated AI extraction and verification of details from uploaded documents (Income, ID, etc.).",
        "• Proactive Assistance: Smart journey tracking alerts users for missing documents or new opportunities.",
        "• Unified Dashboard: A seamless workspace for citizens to track applications and get 24/7 AI assistance."
    ]
    add_content_text(slide, content)

    # 4. ARCHITECTURE
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_template_style(slide, "ARCHITECTURE")
    content = [
        "• Frontend Layer: User-friendly web application for easy citizen access.",
        "• AI Orchestration: Conversational intelligence for natural dialogue and service matching.",
        "• Processing Engine: OCR and NLP pipelines for extracting details from official certificates.",
        "• Data Layer: Secure, compliant storage for user profiles and application tracking."
    ]
    add_content_text(slide, content)

    # 5. TECHNOLOGY USED
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_template_style(slide, "TECHNOLOGY USED")
    content = [
        "• Frontend: React / Next.js, Tailwind CSS",
        "• AI / ML: Large Language Models (LLMs) for reasoning, Computer Vision/OCR for Document Doctor",
        "• Backend & API: Python (FastAPI/Flask) / Node.js",
        "• Database & Storage: PostgreSQL / MongoDB, Secure Cloud Storage"
    ]
    add_content_text(slide, content)

    # 6. WORKING PROTOTYPE
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_template_style(slide, "WORKING PROTOTYPE")
    txBox = slide.shapes.add_textbox(Inches(3), Inches(3.5), Inches(4), Inches(1))
    tf = txBox.text_frame
    tf.text = "Attach video link for demo"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = text_color
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 7. UTILITY/SCALABILITY
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_template_style(slide, "UTILITY/SCALABILITY")
    content = [
        "• Utility: Radically simplifies public service access, bridging the digital divide for millions of citizens.",
        "• Scalability: Cloud-native architecture allows seamless integration of new government schemes.",
        "• Adaptability: Can be easily expanded to support multi-language voice interactions and WhatsApp integration."
    ]
    add_content_text(slide, content)

    # 8. THE JANSAHAY EXPERIENCE (Inspiration from Wellcall)
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_template_style(slide, "THE JANSAHAY EXPERIENCE")
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "(Please insert your 4 screenshots here: Dashboard, Workspace Assistant, Discover Services, Document Doctor)"
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = text_color
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add placeholders for images
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(3.5), Inches(2)).fill.solid()
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3), Inches(3.5), Inches(2)).fill.solid()
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.2), Inches(3.5), Inches(2)).fill.solid()
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(5.2), Inches(3.5), Inches(2)).fill.solid()

    prs.save("JANSAHAY_Presentation.pptx")
    print("Presentation saved as JANSAHAY_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
